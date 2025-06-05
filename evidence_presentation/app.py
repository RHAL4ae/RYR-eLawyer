import os
import logging
from flask import Flask, request, render_template_string, send_from_directory
from pptx import Presentation

app = Flask(__name__)

UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), 'uploads')
OUTPUT_FOLDER = os.path.join(os.path.dirname(__file__), 'outputs')
LOG_FILE = os.path.join(os.path.dirname(__file__), 'presentation.log')

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
logging.basicConfig(filename=LOG_FILE, level=logging.INFO,
                    format='%(asctime)s %(levelname)s %(message)s')

HTML_FORM = """
<!doctype html>
<title>Evidence Uploader</title>
<h1>Upload Evidence (drag & drop or choose files)</h1>
<form method=post enctype=multipart/form-data>
  <input type=file name=files multiple>
  <input type=submit value=Upload>
</form>
"""

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        files = request.files.getlist('files')
        saved = []
        for f in files:
            if not f.filename:
                continue
            dest = os.path.join(UPLOAD_FOLDER, f.filename)
            f.save(dest)
            saved.append(dest)
            logging.info('uploaded %s', f.filename)
        if saved:
            pptx_path = os.path.join(OUTPUT_FOLDER, 'presentation.pptx')
            generate_presentation(saved, pptx_path)
            return send_from_directory(OUTPUT_FOLDER, 'presentation.pptx', as_attachment=True)
    return render_template_string(HTML_FORM)

def generate_presentation(files, pptx_path):
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]
    for path in files:
        filename = os.path.basename(path)
        slide = prs.slides.add_slide(blank_layout)
        title = slide.shapes.title
        if title:
            title.text = filename
        logging.info('added %s to slide', filename)
    prs.save(pptx_path)
    logging.info('presentation saved %s', pptx_path)

if __name__ == '__main__':
    app.run(debug=True)
